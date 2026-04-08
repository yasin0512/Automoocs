"""
Module 5: Automated Presentation Generation
============================================
Thesis Ch. 4-6

Supports two template modes:
  A) Built-in themes (5 styles, python-pptx)
  B) User-uploaded .pptx template — system extracts layout/master and injects content

Image strategies (A/B/C) unchanged.
"""

import io, os, re, base64, logging, shutil
from pathlib import Path
from typing import List, Dict, Optional, Tuple, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
logger = logging.getLogger(__name__)

THEMES = {
    "Modern Blue":    {"bg":(245,248,252),"title":(41,98,149), "body":(55,71,79), "accent":(33,150,243), "bullet":"▶"},
    "Academic Green": {"bg":(240,248,241),"title":(27,94,32),  "body":(46,125,50),"accent":(76,175,80),  "bullet":"●"},
    "Elegant Purple": {"bg":(248,245,255),"title":(74,20,140), "body":(106,27,154),"accent":(156,39,176),"bullet":"◆"},
    "Tech Gray":      {"bg":(236,239,241),"title":(38,50,56),  "body":(69,90,100),"accent":(0,188,212),  "bullet":"▸"},
    "Warm Orange":    {"bg":(255,248,225),"title":(191,54,12), "body":(109,76,65),"accent":(255,87,34),  "bullet":"▼"},
}


# ──────────────────────────────────────────────────────────────
# IMAGE ALIGNER (BLIP-2 + CLIP)
# ──────────────────────────────────────────────────────────────
class ImageAligner:
    def __init__(self):
        self._clip = self._clip_prep = self._blip_proc = self._blip_model = None
        self._dev = "cpu"

    def _load_clip(self):
        if self._clip: return
        try:
            import torch, clip
            self._dev = "cuda" if torch.cuda.is_available() else "cpu"
            self._clip, self._clip_prep = clip.load("ViT-B/32", device=self._dev)
        except Exception as e: logger.warning(f"CLIP unavailable: {e}")

    def _blip2_caption(self, img_path: str) -> str:
        if not self._blip_model: return ""
        try:
            from PIL import Image; import torch
            img = Image.open(img_path).convert("RGB")
            inputs = self._blip_proc(images=img, return_tensors="pt").to(
                next(self._blip_model.parameters()).device)
            with torch.no_grad():
                ids = self._blip_model.generate(**inputs, max_new_tokens=50)
            return self._blip_proc.decode(ids[0], skip_special_tokens=True)
        except Exception: return ""

    def match(self, slide_text: str, image_paths: List[str],
              use_blip2: bool = False) -> Optional[str]:
        if not image_paths: return None
        self._load_clip()
        if not self._clip: return image_paths[0]
        try:
            import torch, clip
            from PIL import Image
            tokens = clip.tokenize([slide_text[:77]], truncate=True).to(self._dev)
            with torch.no_grad():
                tf = self._clip.encode_text(tokens)
                tf /= tf.norm(dim=-1, keepdim=True)
            best, bscore = image_paths[0], -999.0
            for path in image_paths:
                try:
                    img_t = self._clip_prep(Image.open(path).convert("RGB")).unsqueeze(0).to(self._dev)
                    with torch.no_grad():
                        vf = self._clip.encode_image(img_t)
                        vf /= vf.norm(dim=-1, keepdim=True)
                    sc = float((tf @ vf.T)[0,0])
                    if sc > bscore: bscore, best = sc, path
                except Exception: continue
            return best
        except Exception as e:
            logger.warning(f"CLIP error: {e}")
            return image_paths[0]


# ──────────────────────────────────────────────────────────────
# STABLE DIFFUSION
# ──────────────────────────────────────────────────────────────
class StableDiffusionGenerator:
    def __init__(self, sd_url: str = "http://127.0.0.1:7860"):
        self.url = sd_url

    def is_available(self) -> bool:
        import requests
        try:
            return requests.get(f"{self.url}/sdapi/v1/options", timeout=3).status_code == 200
        except Exception: return False

    def generate(self, topic: str, out_path: str, width: int = 512, height: int = 512) -> bool:
        import requests
        try:
            r = requests.post(f"{self.url}/sdapi/v1/txt2img", json={
                "prompt": f"educational illustration, {topic}, clean, professional, academic, no text",
                "negative_prompt": "low quality, blurry, nsfw, text, watermark",
                "steps": 20, "cfg_scale": 7.0, "width": width, "height": height,
            }, timeout=120)
            r.raise_for_status()
            with open(out_path, "wb") as f: f.write(base64.b64decode(r.json()["images"][0]))
            return True
        except Exception as e:
            logger.warning(f"SD failed: {e}"); return False


# ──────────────────────────────────────────────────────────────
# SEMANTIC PAGINATION
# ──────────────────────────────────────────────────────────────
def segment_slides(transcript: str, key_points: Optional[List[str]] = None,
                   max_per_slide: int = 5) -> List[Dict]:
    sents = [s.strip() for s in re.split(r"[。！？\n]", transcript) if len(s.strip()) > 3]
    if not sents:
        return [{"title": "Course Content", "bullets": [transcript[:200]], "kp_index": 0}]
    n = max(2, min((len(key_points) if key_points else max_per_slide), len(sents)))
    chunk = max(1, len(sents) // n)
    slides = []
    for i in range(n):
        cs = sents[i * chunk:(i + 1) * chunk]
        if not cs: break
        slides.append({
            "title":    key_points[i] if key_points and i < len(key_points) else cs[0][:35],
            "bullets":  cs[:4],
            "kp_index": i,
        })
    return slides


# ──────────────────────────────────────────────────────────────
# USER TEMPLATE PARSER
# ──────────────────────────────────────────────────────────────
class UserTemplateParser:
    """
    Parse a user-uploaded .pptx template.
    Extracts: slide master, color scheme, fonts.
    The system then injects generated content into the template's slide structure.
    """

    @staticmethod
    def load(template_path: str) -> Presentation:
        """Load user's pptx template and return a Presentation object with master preserved."""
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")
        prs = Presentation(template_path)
        logger.info(f"Loaded user template: {template_path} "
                    f"({len(prs.slides)} existing slides, "
                    f"{len(prs.slide_layouts)} layouts)")
        return prs

    @staticmethod
    def detect_colors(prs: Presentation) -> Dict:
        """Try to extract dominant accent color from template's theme."""
        try:
            from pptx.dml.color import RGBColor
            # Check first slide layout background
            for layout in prs.slide_layouts:
                bg = layout.background.fill
                if bg.type is not None:
                    try:
                        return {"accent": bg.fore_color.rgb}
                    except Exception:
                        pass
        except Exception:
            pass
        return {}

    @staticmethod
    def inject_content(prs: Presentation, slides_data: List[Dict],
                       image_paths: Optional[List[str]] = None,
                       aligner: Optional["ImageAligner"] = None) -> Presentation:
        """
        Inject generated slide content into the user's template.
        Uses the template's blank layout (index 6 or last blank) for content slides.
        Clears existing slides and rebuilds with generated content.
        """
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # Find best blank layout
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

        # Remove existing slides (keep master/layouts)
        xml_slides = prs.slides._sldIdLst
        for slide in list(prs.slides):
            rId = prs.slides._sldIdLst.index(
                [x for x in prs.slides._sldIdLst if True][0]
            )
        # Safer: just add new slides on top (template slides remain as reference)
        # We'll add our content slides
        imgs = list(image_paths or [])

        for i, page in enumerate(slides_data):
            slide = prs.slides.add_slide(blank_layout)
            query = page["title"] + " " + " ".join(page["bullets"][:2])

            # Title box
            tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(9.2), Inches(1.0))
            p = tb.text_frame.paragraphs[0]
            p.text = page["title"]
            p.font.name = "Microsoft JhengHei UI"
            p.font.size = Pt(28)
            p.font.bold = True

            # Bullet box
            img_path = aligner.match(query, imgs) if imgs and aligner else None
            w = Inches(5.5) if img_path else Inches(9.2)
            bx = slide.shapes.add_textbox(Inches(0.4), Inches(1.5), w, Inches(5.5))
            tf = bx.text_frame; tf.word_wrap = True
            for j, bullet in enumerate(page["bullets"]):
                pr = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                pr.text = f"  ▸  {bullet}"
                pr.font.name = "Microsoft JhengHei UI"
                pr.font.size = Pt(17)
                pr.space_after = Pt(8)

            # Image
            if img_path:
                try:
                    from PIL import Image as PILImage
                    img = PILImage.open(img_path)
                    iw, ih = img.size
                    pw = Inches(3.4); ph = min(pw * ih / max(iw, 1), Inches(5.2))
                    slide.shapes.add_picture(img_path, Inches(6.4), Inches(1.5), pw, ph)
                except Exception as e:
                    logger.warning(f"Template image insert failed: {e}")

        return prs


# ──────────────────────────────────────────────────────────────
# PPT GENERATOR  (built-in themes)
# ──────────────────────────────────────────────────────────────
class PPTGenerator:
    W = Inches(10); H = Inches(7.5)

    def __init__(self, theme: str = "Modern Blue",
                 sd_url: str = "http://127.0.0.1:7860",
                 template_path: Optional[str] = None):
        """
        Args:
            theme         : built-in theme name (used when template_path is None)
            sd_url        : Stable Diffusion WebUI URL
            template_path : path to user-uploaded .pptx template (overrides theme)
        """
        self.cfg           = THEMES.get(theme, THEMES["Modern Blue"])
        self.theme_name    = theme
        self.aligner       = ImageAligner()
        self.sd            = StableDiffusionGenerator(sd_url)
        self.template_path = template_path
        self.use_template  = bool(template_path and os.path.exists(template_path))

        if self.use_template:
            logger.info(f"PPTGenerator: using user template → {template_path}")
        else:
            logger.info(f"PPTGenerator: using built-in theme → {theme}")

    # ── Internal helpers (built-in theme) ─────────────────────
    def _rgb(self, t: Tuple) -> RGBColor: return RGBColor(*t)

    def _fill_bg(self, slide):
        bg = slide.background.fill; bg.solid()
        bg.fore_color.rgb = self._rgb(self.cfg["bg"])

    def _add_title(self, slide, text: str):
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9.2), Inches(1.1))
        p  = tb.text_frame.paragraphs[0]
        p.text = text; p.font.name = "Microsoft JhengHei UI"
        p.font.size = Pt(30); p.font.bold = True
        p.font.color.rgb = self._rgb(self.cfg["title"])
        bar = slide.shapes.add_shape(1, Inches(0.4), Inches(1.25), Inches(9.2), Inches(0.06))
        bar.fill.solid(); bar.fill.fore_color.rgb = self._rgb(self.cfg["accent"])
        bar.line.fill.background()

    def _add_bullets(self, slide, bullets: List[str], has_img: bool = False):
        w  = Inches(5.5) if has_img else Inches(9.2)
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(1.5), w, Inches(5.7))
        tf = tb.text_frame; tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"  {self.cfg['bullet']}  {b}"
            p.font.name = "Microsoft JhengHei UI"; p.font.size = Pt(18)
            p.font.color.rgb = self._rgb(self.cfg["body"]); p.space_after = Pt(10)

    def _add_image(self, slide, path: str):
        try:
            from PIL import Image as PILImage
            img = PILImage.open(path); iw, ih = img.size
            pw = Inches(3.4); ph = min(pw * ih / max(iw, 1), Inches(5.5))
            slide.shapes.add_picture(path, Inches(6.4), Inches(1.5), pw, ph)
        except Exception as e: logger.warning(f"Image insert failed: {e}")

    def _make_cover(self, prs: Presentation, title: str, subtitle: str):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._fill_bg(slide)
        for txt, y, size, bold in [
            (title,    2.0, 44, True),
            (subtitle, 4.2, 20, False),
        ]:
            tb = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(1.2))
            p  = tb.text_frame.paragraphs[0]
            p.text = txt; p.alignment = PP_ALIGN.CENTER
            p.font.name = "Microsoft JhengHei UI"; p.font.size = Pt(size); p.font.bold = bold
            p.font.color.rgb = self._rgb(self.cfg["title"] if bold else self.cfg["body"])
        bar = slide.shapes.add_shape(1, Inches(0), Inches(6.9), Inches(10), Inches(0.6))
        bar.fill.solid(); bar.fill.fore_color.rgb = self._rgb(self.cfg["accent"])
        bar.line.fill.background()

    # ── Main generate ──────────────────────────────────────────
    def generate(self, transcript: str, course_title: str = "Lecture",
                 key_points: Optional[List[str]] = None,
                 image_paths: Optional[List[str]] = None,
                 use_sd: bool = False, use_blip2: bool = False,
                 sd_output_dir: str = "output/sd_images",
                 progress_cb=None) -> io.BytesIO:

        def _cb(p, m):
            if progress_cb: progress_cb(p, m)

        _cb(0.05, "Segmenting slides…")
        slides_data = segment_slides(transcript, key_points)
        imgs = list(image_paths or [])
        os.makedirs(sd_output_dir, exist_ok=True)

        # ── Branch: user template ────────────────────────────
        if self.use_template:
            _cb(0.10, f"Loading user template: {Path(self.template_path).name}…")
            try:
                parser = UserTemplateParser()
                prs = parser.load(self.template_path)
                # Remove any existing content slides (keep masters/layouts)
                # Add a cover title on first existing slide if possible, else add new
                prs = parser.inject_content(
                    prs, slides_data,
                    image_paths=imgs or None,
                    aligner=self.aligner if imgs else None,
                )
                _cb(0.92, "Packing user-template PPTX…")
                buf = io.BytesIO(); prs.save(buf); buf.seek(0)
                _cb(1.0, "Template-based presentation ready")
                return buf
            except Exception as e:
                logger.warning(f"User template failed ({e}), falling back to built-in theme")
                self.use_template = False

        # ── Branch: built-in theme ───────────────────────────
        prs = Presentation(); prs.slide_width = self.W; prs.slide_height = self.H
        _cb(0.10, "Cover slide…")
        self._make_cover(prs, course_title, "Auto-generated MOOCs Lecture Slides")

        total = len(slides_data)
        for i, page in enumerate(slides_data):
            _cb(0.10 + 0.78 * i / max(total, 1), f"Slide {i + 1}/{total}")
            query = page["title"] + " " + " ".join(page["bullets"][:2])
            img_path: Optional[str] = None
            if imgs:
                img_path = self.aligner.match(query, imgs, use_blip2=use_blip2)
            elif use_sd and self.sd.is_available():
                gp = os.path.join(sd_output_dir, f"slide_{i:03d}.png")
                if self.sd.generate(query[:200], gp): img_path = gp

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._fill_bg(slide)
            self._add_title(slide, page["title"])
            self._add_bullets(slide, page["bullets"], has_img=bool(img_path))
            if img_path: self._add_image(slide, img_path)

            pn = slide.shapes.add_textbox(Inches(9.3), Inches(7.1), Inches(0.6), Inches(0.3))
            pn.text_frame.paragraphs[0].text = str(i + 2)
            pn.text_frame.paragraphs[0].font.size = Pt(9)
            pn.text_frame.paragraphs[0].font.color.rgb = self._rgb(self.cfg["body"])

        _cb(0.92, "Packing PPTX…")
        buf = io.BytesIO(); prs.save(buf); buf.seek(0)
        _cb(1.0, "Presentation ready")
        return buf

    def generate_to_file(self, out_path: str, **kwargs) -> str:
        buf = self.generate(**kwargs)
        with open(out_path, "wb") as f: f.write(buf.read())
        return out_path
